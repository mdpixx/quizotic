#!/usr/bin/env python3
"""Parse a PPTX file and output structured JSON to stdout.

Usage: python3 parse_pptx.py <path-to-pptx>

Output: JSON array of slide objects with structure, hierarchy, and metadata.
"""

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def classify_slide(slide, layout_name, placeholders, shapes):
    """Determine slide type from layout name and content."""
    layout_lower = layout_name.lower()

    # Count content types
    has_table = any(
        shape.has_table for shape in shapes
    )
    has_chart = any(
        shape.has_chart for shape in shapes if hasattr(shape, 'has_chart')
    )
    has_smartart = any(
        shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in shapes
    )
    image_count = sum(
        1 for shape in shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        or (hasattr(shape, 'image') and shape.shape_type != MSO_SHAPE_TYPE.PLACEHOLDER)
    )
    text_length = sum(
        len(shape.text_frame.text) for shape in shapes
        if shape.has_text_frame
    )

    if has_table:
        return 'table'

    # Title slide: layout says "title" and has few placeholders
    if ('title' in layout_lower and 'content' not in layout_lower) and len(placeholders) <= 3:
        return 'title'

    if 'section' in layout_lower:
        return 'section'

    if 'blank' in layout_lower and text_length < 10 and image_count == 0:
        return 'blank'

    # Image-dominant: has images but minimal text
    if image_count > 0 and text_length < 50:
        return 'image_dominant'

    # No content at all
    if text_length < 5 and image_count == 0 and not has_table and not has_chart:
        return 'blank'

    return 'content'


def extract_paragraphs(text_frame):
    """Extract paragraphs with hierarchy info from a text frame."""
    paragraphs = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        # Determine formatting from runs
        is_bold = False
        font_size = None
        for run in para.runs:
            if run.font.bold:
                is_bold = True
            if run.font.size and font_size is None:
                font_size = run.font.size

        paragraphs.append({
            'text': text,
            'level': para.level or 0,
            'bold': is_bold,
            'fontSize': round(font_size / Pt(1), 1) if font_size else None,
        })
    return paragraphs


def extract_table(shape):
    """Extract table data as a 2D array of strings."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text.strip())
        rows.append(cells)
    return rows


def extract_notes(slide):
    """Extract speaker notes text."""
    if not slide.has_notes_slide:
        return None
    notes_frame = slide.notes_slide.notes_text_frame
    text = notes_frame.text.strip()
    if not text or len(text) < 3:
        return None
    # Filter out boilerplate (just a number, or "slide N")
    import re
    if re.match(r'^\d+$', text) or re.match(r'^slide\s*\d*$', text, re.IGNORECASE):
        return None
    return text


def parse_pptx(filepath):
    """Parse PPTX and return structured slide data."""
    prs = Presentation(filepath)
    slides_data = []

    for idx, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name if slide.slide_layout else 'Unknown'
        shapes = list(slide.shapes)
        placeholders = list(slide.placeholders)

        # Classify slide type
        slide_type = classify_slide(slide, layout_name, placeholders, shapes)

        # Extract title and subtitle from placeholders
        title = None
        subtitle = None
        body_paragraphs = []

        for ph in placeholders:
            if not ph.has_text_frame:
                continue
            ph_type = ph.placeholder_format.type if ph.placeholder_format else None
            ph_idx = ph.placeholder_format.idx if ph.placeholder_format else None

            # Title placeholder (type TITLE=15, CENTER_TITLE=3)
            if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE) or ph_idx == 0:
                title_text = ph.text_frame.text.strip()
                if title_text:
                    title = title_text
            # Subtitle placeholder (type SUBTITLE=4)
            elif ph_type == PP_PLACEHOLDER.SUBTITLE or ph_idx == 1:
                if slide_type in ('title', 'section'):
                    subtitle_text = ph.text_frame.text.strip()
                    if subtitle_text:
                        subtitle = subtitle_text
                else:
                    # Body content
                    body_paragraphs.extend(extract_paragraphs(ph.text_frame))
            # Body placeholder (type BODY=2, OBJECT=7)
            elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT) or ph_idx >= 2:
                body_paragraphs.extend(extract_paragraphs(ph.text_frame))

        # Also extract text from non-placeholder shapes
        for shape in shapes:
            if shape in placeholders:
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) > 5:
                    body_paragraphs.extend(extract_paragraphs(shape.text_frame))

        # Extract tables
        table_data = None
        for shape in shapes:
            if shape.has_table:
                table_data = extract_table(shape)
                break  # Take first table

        # Count images
        image_count = sum(
            1 for shape in shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        )

        # Charts
        has_chart = any(
            shape.has_chart for shape in shapes if hasattr(shape, 'has_chart')
        )

        # SmartArt (approximated as grouped shapes)
        has_smartart = any(
            shape.shape_type == MSO_SHAPE_TYPE.GROUP for shape in shapes
        )

        # Speaker notes
        notes = extract_notes(slide)

        # If no title found from placeholders, use first substantial paragraph
        if not title and body_paragraphs:
            title = body_paragraphs[0]['text']
            body_paragraphs = body_paragraphs[1:]

        slides_data.append({
            'index': idx,
            'slideType': slide_type,
            'title': title,
            'subtitle': subtitle,
            'bodyParagraphs': body_paragraphs,
            'tableData': table_data,
            'speakerNotes': notes,
            'imageCount': image_count,
            'hasChart': has_chart,
            'hasSmartArt': has_smartart,
            'layoutName': layout_name,
        })

    return slides_data


if __name__ == '__main__':
    if len(sys.argv) != 2:
        print(json.dumps({'error': 'Usage: parse_pptx.py <filepath>'}), file=sys.stderr)
        sys.exit(1)

    filepath = sys.argv[1]
    if not Path(filepath).exists():
        print(json.dumps({'error': f'File not found: {filepath}'}), file=sys.stderr)
        sys.exit(1)

    try:
        result = parse_pptx(filepath)
        print(json.dumps(result, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({'error': str(e)}), file=sys.stderr)
        sys.exit(1)
